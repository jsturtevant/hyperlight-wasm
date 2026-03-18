"""Basic example of hyperlight_sandbox Python SDK."""

import asyncio
import time
from pathlib import Path
from hyperlight_sandbox import WasmSandbox

def timed_run(sandbox, code, label="run"):
    """Run code in sandbox and print timing."""
    start = time.perf_counter()
    result = sandbox.run(code)
    elapsed_ms = (time.perf_counter() - start) * 1000
    print(f"⏱️  {label}: {elapsed_ms:.1f}ms")
    return result

# Create sandbox pointing to the AOT-compiled Python component
t0 = time.perf_counter()
sandbox = WasmSandbox(module_path="src/python_sandbox/python-sandbox.aot")
print(f"⏱️  WasmSandbox created (lazy): {(time.perf_counter() - t0) * 1000:.1f}ms")

# Register host tools before first run()
sandbox.register_tool("add", lambda a=0, b=0: a + b)
sandbox.register_tool("greet", lambda name="world": f"Hello, {name}!")

# Async functions work too — no wrapping needed
async def async_multiply(a=0, b=0):
    await asyncio.sleep(0.5)  # simulate async I/O
    return a * b

sandbox.register_tool("multiply", async_multiply)

# Host tool using python-pptx (C extension — can't run in Wasm guest)
def make_pptx(title="", subtitle="", slides=None):
    """Build a PPTX on the host side using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = subtitle

    # Content slides
    for slide_data in (slides or []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data.get("title", "")
        slide.placeholders[1].text = slide_data.get("body", "")

    prs.save("demo.pptx")
    return f"{len(slides or [])} slides saved to demo.pptx"

sandbox.register_tool("make_pptx", make_pptx)

# Test 1: Basic code execution (first run triggers sandbox init)
print("\n--- Test 1: Basic execution (includes sandbox init) ---")
result = timed_run(sandbox, 'print("hello from python SDK!")', "first run (cold)")
print(f"stdout: {result.stdout!r}")
print(f"success: {result.success}")

# Test 2: Tool dispatch via call_tool() — sync and async
print("\n--- Test 2: Tool dispatch (sync + async) ---")
result = timed_run(sandbox, """
import time
result = call_tool('add', a=3, b=4)
greeting = call_tool('greet', name='James')
t0 = time.time()
product = call_tool('multiply', a=6, b=7)
elapsed = time.time() - t0
print(f"3 + 4 = {result}")
print(f"{greeting}")
print(f"6 * 7 = {product}  (async tool, slept {elapsed:.1f}s)")
try:
    call_tool('nonexistent', x=1)
except RuntimeError as e:
    print(f"Caught error: {e}")
print("All tool tests passed!")
""", "tool dispatch")
print(f"stdout: {result.stdout!r}")
print(f"success: {result.success}")

# Test 3: Sandbox reuse
print("\n--- Test 3: Sandbox reuse ---")
result = timed_run(sandbox, 'print("second run works!")', "reuse (warm)")
print(f"stdout: {result.stdout!r}")
print(f"success: {result.success}")

# Test 4: Snapshot/restore
print("\n--- Test 4: Snapshot/restore ---")
t0 = time.perf_counter()
snap = sandbox.snapshot()
print(f"⏱️  snapshot: {(time.perf_counter() - t0) * 1000:.1f}ms")

result1 = timed_run(sandbox, 'x = 42; print(f"x = {x}")', "pre-restore run")
print(f"Before restore: {result1.stdout!r}")

t0 = time.perf_counter()
sandbox.restore(snap)
print(f"⏱️  restore: {(time.perf_counter() - t0) * 1000:.1f}ms")

result2 = timed_run(sandbox, """
try:
    print(f"x = {x}")
except NameError:
    print("x is not defined (state was rolled back)")
""", "post-restore run")
print(f"After restore: {result2.stdout!r}")

# Test 5: File I/O via WASI filesystem
print("\n--- Test 5: File I/O ---")
sandbox.add_file("data.json", b'{"greeting": "hello from file!"}')
result = timed_run(sandbox, """
import json
with open('/input/data.json', 'r') as f:
    data = json.load(f)
print(f"Read: {data['greeting']}")

with open('/output/result.txt', 'w') as f:
    f.write('Written from Python SDK!')
print("File I/O works!")
""", "file I/O")
print(f"stdout: {result.stdout!r}")
print(f"outputs: {dict((k, v.decode()) for k, v in result.outputs.items())}")
assert result.success
assert result.outputs["result.txt"] == b"Written from Python SDK!"

# Test 6: Generate a PowerPoint file via host tool
# python-pptx requires lxml (C extension) which can't run in Wasm,
# so the guest prepares slide data and the host builds the PPTX.
print("\n--- Test 6: PowerPoint generation (guest data → host pptx) ---")
result = timed_run(sandbox, """
import json

slides = [
    {"title": "Hyperlight Sandbox", "body": "Code runs in hardware-isolated Wasm components"},
    {"title": "Key Features", "body": "Tool dispatch, File I/O, Snapshots, WASI-HTTP networking"},
    {"title": "Guest + Host", "body": "Guest computes data, host builds the PPTX with python-pptx"},
]
result = call_tool('make_pptx', title='Hyperlight Demo', subtitle='Generated from a Wasm sandbox', slides=slides)
print(f"PPTX created: {result}")
""", "pptx via host tool")
print(f"stdout: {result.stdout!r}")
assert result.success
pptx_path = Path("demo.pptx")
assert pptx_path.exists(), "demo.pptx not written"
size = pptx_path.stat().st_size
print(f"Host wrote: demo.pptx ({size:,} bytes)")
# Verify it's a valid PPTX
import zipfile
with zipfile.ZipFile(pptx_path) as zf:
    slide_files = [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
    print(f"PPTX contains: {len(slide_files)} slides — {slide_files}")
assert len(slide_files) == 4  # title + 3 content slides
pptx_path.unlink()

print("\n✅ All tests passed!")
